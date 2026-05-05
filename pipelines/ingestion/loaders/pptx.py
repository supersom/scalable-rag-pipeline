# pipelines/ingestion/loaders/pptx.py
import io
from pptx import Presentation
from pptx.util import Pt


def parse_pptx_bytes(file_bytes: bytes, filename: str):
    """
    Extracts text from a PowerPoint file.
    Pulls slide titles, body text, and table cell content in slide order.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)
        if parts:
            slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(parts))

    full_text = "\n\n".join(slides_text)
    metadata = {
        "filename": filename,
        "type": "pptx",
        "slide_count": len(prs.slides),
    }
    return full_text, metadata
